"""
Lambda 1 — Document Ingestion.

Triggered by S3 ObjectCreated on the raw bucket.
Extracts text + markers from PDF/PPTX/TXT, writes JSON to silver bucket.

Mutation 2 (Auto-Scaling Lambda with DLQ): if extraction fails 3 times,
SQS/S3 retries exhaust and the event routes to the DLQ.
"""

import hashlib
import json
import os
import re
import urllib.parse
from datetime import datetime, timezone

import boto3

RAW_BUCKET = os.environ["RAW_BUCKET"]
SILVER_BUCKET = os.environ["SILVER_BUCKET"]
AWS_REGION = os.environ.get("AWS_REGION", "eu-north-1")

s3 = boto3.client("s3", region_name=AWS_REGION)


# ── Marker extraction ──────────────────────────────────────────────────────────

DATE_RE = re.compile(
    r'\b(\d{1,2}[/-]\d{1,2}[/-]\d{2,4}|\d{4}-\d{2}-\d{2}'
    r'|(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\.?\s+\d{1,2},?\s+\d{4})\b',
    re.IGNORECASE,
)
METRIC_RE = re.compile(
    r'\b\d+(?:\.\d+)?\s*(?:ms|s|sec|%|MB|GB|TB|KB|USD|\$|req/s|rpm|RPS|vCPU|GiB)\b',
    re.IGNORECASE,
)
ERROR_RE = re.compile(
    r'\b(?:ERR(?:OR)?[-_]?\d+|[45]\d{2}\s+(?:error|timeout)|exception|fatal)\b',
    re.IGNORECASE,
)


def extract_markers(text):
    return {
        "dates": list(dict.fromkeys(DATE_RE.findall(text))),
        "metrics": list(dict.fromkeys(METRIC_RE.findall(text))),
        "error_codes": list(dict.fromkeys(ERROR_RE.findall(text))),
    }


# ── Text extraction ────────────────────────────────────────────────────────────

def extract_pdf(content: bytes) -> str:
    import io
    import PyPDF2
    reader = PyPDF2.PdfReader(io.BytesIO(content))
    return "\n".join(p.extract_text() or "" for p in reader.pages)


def extract_pptx(content: bytes) -> str:
    import io
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def extract_txt(content: bytes) -> str:
    return content.decode("utf-8", errors="replace")


EXTRACTORS = {".pdf": extract_pdf, ".pptx": extract_pptx, ".txt": extract_txt, ".md": extract_txt}


def doc_id(key: str) -> str:
    return hashlib.sha256(key.encode()).hexdigest()[:16]


# ── Lambda handler ─────────────────────────────────────────────────────────────

def handler(event, context):
    for record in event.get("Records", []):
        bucket = record["s3"]["bucket"]["name"]
        key = urllib.parse.unquote_plus(record["s3"]["object"]["key"])
        suffix = "." + key.rsplit(".", 1)[-1].lower() if "." in key else ""

        extractor = EXTRACTORS.get(suffix)
        if not extractor:
            print(f"Unsupported file type: {key}")
            continue

        print(f"Extracting: s3://{bucket}/{key}")
        obj = s3.get_object(Bucket=bucket, Key=key)
        content = obj["Body"].read()
        text = extractor(content)
        markers = extract_markers(text)

        record_out = {
            "doc_id": doc_id(key),
            "file_name": key.split("/")[-1],
            "file_type": suffix.lstrip("."),
            "version": 1,
            "uploaded_at": datetime.now(timezone.utc).isoformat(),
            "char_count": len(text),
            "markers": markers,
            "text": text,
        }

        silver_key = f"silver/{record_out['doc_id']}_v1.json"
        s3.put_object(
            Bucket=SILVER_BUCKET,
            Key=silver_key,
            Body=json.dumps(record_out).encode(),
            ContentType="application/json",
        )
        print(f"Written to s3://{SILVER_BUCKET}/{silver_key}")

    return {"statusCode": 200}
