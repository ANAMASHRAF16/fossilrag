"""
Component 1 — Serverless Document Ingestion (local runner).

Accepts PDF, PPTX, or TXT files from RAW_DIR.
Extracts text + metadata + markers, writes JSON to SILVER_DIR.

In production this runs as a Lambda triggered by S3 ObjectCreated.
Locally we run it as a script: python -m src.extractor
"""

import hashlib
import json
import os
import re
from datetime import datetime, timezone
from pathlib import Path

from dotenv import load_dotenv

load_dotenv()

RAW_DIR = Path(os.environ.get("RAW_DIR", "data/raw"))
SILVER_DIR = Path(os.environ.get("SILVER_DIR", "data/silver"))
SILVER_DIR.mkdir(parents=True, exist_ok=True)


def extract_pdf(path: Path) -> str:
    import PyPDF2
    text_parts = []
    with open(path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            text_parts.append(page.extract_text() or "")
    return "\n".join(text_parts)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def extract_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


EXTRACTORS = {
    ".pdf": extract_pdf,
    ".pptx": extract_pptx,
    ".txt": extract_txt,
    ".md": extract_txt,
}


def doc_id(path: Path) -> str:
    return hashlib.sha256(path.name.encode()).hexdigest()[:16]


def process_file(path: Path, version: int = 1) -> dict:
    suffix = path.suffix.lower()
    extractor = EXTRACTORS.get(suffix)
    if not extractor:
        print(f"  Skipping unsupported file: {path.name}")
        return {}

    print(f"  Extracting: {path.name}")
    text = extractor(path)

    from src.markers import extract_markers
    markers = extract_markers(text)

    record = {
        "doc_id": doc_id(path),
        "file_name": path.name,
        "file_type": suffix.lstrip("."),
        "version": version,
        "uploaded_at": datetime.now(timezone.utc).isoformat(),
        "char_count": len(text),
        "markers": markers,
        "text": text,
    }

    out_path = SILVER_DIR / f"{record['doc_id']}_v{version}.json"
    out_path.write_text(json.dumps(record, indent=2), encoding="utf-8")
    print(f"    -> {out_path.name} ({len(text)} chars, markers: {sum(len(v) for v in markers.values())} found)")
    return record


def run(version: int = 1):
    files = [f for f in RAW_DIR.iterdir() if f.is_file() and f.suffix.lower() in EXTRACTORS]
    if not files:
        print(f"No supported files in {RAW_DIR}. Drop a PDF, PPTX, or TXT file there and re-run.")
        return
    print(f"Extracting {len(files)} file(s) from {RAW_DIR}...")
    for f in files:
        process_file(f, version=version)
    print("Extraction complete.")


if __name__ == "__main__":
    import sys
    version = int(sys.argv[1]) if len(sys.argv) > 1 else 1
    run(version=version)
